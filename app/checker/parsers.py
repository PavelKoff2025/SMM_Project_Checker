
import pdfplumber
from pptx import Presentation
from PyPDF2 import PdfReader


class PDFParser:

    @staticmethod
    def extract_text(file_path: str) -> dict:
        pages = []
        full_text = []

        with pdfplumber.open(file_path) as pdf:
            for i, page in enumerate(pdf.pages, 1):
                text = page.extract_text() or ''
                pages.append({'number': i, 'text': text})
                full_text.append(text)

        metadata = {}
        try:
            reader = PdfReader(file_path)
            meta = reader.metadata
            if meta:
                metadata = {
                    'title': meta.get('/Title'),
                    'author': meta.get('/Author'),
                    'subject': meta.get('/Subject'),
                    'pages': len(reader.pages),
                }
        except Exception:
            metadata = {'pages': len(pages)}

        return {
            'text': '\n'.join(full_text).strip(),
            'pages': pages,
            'metadata': metadata,
        }

    @staticmethod
    def ocr_fallback(file_path: str) -> str | None:
        try:
            import pytesseract
            from pdf2image import convert_from_path
        except ImportError:
            return None

        try:
            images = convert_from_path(file_path, dpi=300)
            text_parts = []
            for img in images:
                text = pytesseract.image_to_string(img, lang='rus+eng')
                text_parts.append(text)
            return '\n'.join(text_parts).strip()
        except Exception:
            return None

    @classmethod
    def extract_with_fallback(cls, file_path: str) -> dict:
        result = cls.extract_text(file_path)
        if not result['text'].strip():
            ocr_text = cls.ocr_fallback(file_path)
            if ocr_text:
                result['text'] = ocr_text
                result['ocr_used'] = True
        return result


class PPTXParser:

    @staticmethod
    def extract_text(file_path: str) -> dict:
        prs = Presentation(file_path)
        slides = []
        full_text = []

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_title = ''
            slide_content = []

            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue

                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if not text:
                        continue

                    if shape == slide.shapes.title:
                        slide_title = text
                    else:
                        slide_content.append(text)

            slides.append({
                'number': slide_num,
                'title': slide_title,
                'content': slide_content,
            })
            full_text.append(slide_title)
            full_text.extend(slide_content)

        return {
            'text': '\n'.join(full_text).strip(),
            'slides': slides,
        }
